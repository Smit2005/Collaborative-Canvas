import os
import json
import re
import sys
import unicodedata
from pptx import Presentation
from transformers import pipeline

# Force UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

def clean_text(text):
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\xa0", " ")
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_text_and_images(pptx_path, output_dir="output"):
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    slides = []

    for i, slide in enumerate(prs.slides):
        text_list = []
        image_path = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                cleaned = clean_text(shape.text)
                if cleaned:
                    text_list.append(cleaned)

            if shape.shape_type == 13:  # Picture
                image = shape.image
                ext = image.ext
                img_name = f"slide_{i+1}_img.{ext}"
                img_path = os.path.join(output_dir, img_name)
                with open(img_path, "wb") as f:
                    f.write(image.blob)
                image_path = img_name

        full_text = "\n".join(text_list)
        slides.append({
            "slide_number": i + 1,
            "text": full_text,
            "word_count": len(full_text.split()),
            "image_path": image_path
        })

    return slides

def create_batches(slides, min_words=100):
    batches = []
    batch = []
    word_total = 0

    for slide in slides:
        batch.append(slide)
        word_total += slide["word_count"]

        if word_total >= min_words:
            batches.append(batch)
            batch = []
            word_total = 0

    if batch:
        batches.append(batch)

    return batches

def summarize_batches(batches):
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    results = []

    for batch in batches:
        combined_text = "\n".join(s["text"] for s in batch if s["text"])
        if len(combined_text.split()) < 50:
            continue

        summary = summarizer(combined_text, max_length=300, min_length=100, do_sample=False)[0]["summary_text"]
        results.append({
            "slide": [batch[0]["slide_number"], batch[-1]["slide_number"]],
            "summary": summary.strip(),
            "image": [s["image_path"] for s in batch if s["image_path"]]
        })

    return results

# ✅ RENAMED THIS FUNCTION FROM 'main' TO 'run_pipeline'
def run_pipeline(pptx_path, output_dir="output_images"):
    slides = extract_text_and_images(pptx_path, output_dir=output_dir)
    batches = create_batches(slides)
    summaries = summarize_batches(batches)
    return summaries

# This part is for running the file directly from the command line, it can be removed or kept.
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
    else:
        pptx_path_arg = sys.argv[1]
        output = run_pipeline(pptx_path_arg)
        print(json.dumps(output, ensure_ascii=False))